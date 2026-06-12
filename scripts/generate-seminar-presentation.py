#!/usr/bin/env python3
"""Generate Yoqubkhoja Hub seminar presentation as .pptx and .docx."""

from __future__ import annotations

import sys
from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptxRGB
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptxInches, Pt as PptxPt

SCRIPT_DIR = Path(__file__).resolve().parent
sys.path.insert(0, str(SCRIPT_DIR))
from generate_presentation_assets import build_assets  # noqa: E402

OUTPUT_DIR = Path(__file__).resolve().parent.parent / "docs" / "presentation"
PPTX_PATH = OUTPUT_DIR / "Yoqubkhoja-Hub-Seminar.pptx"
DOCX_PATH = OUTPUT_DIR / "Yoqubkhoja-Hub-Seminar.docx"

# Brand palette (src/app/globals.css)
BLUE = PptxRGB(0x3B, 0x82, 0xF6)
BLUE_DARK = PptxRGB(0x25, 0x63, 0xEB)
INDIGO = PptxRGB(0x63, 0x66, 0xF1)
PURPLE = PptxRGB(0x8B, 0x5C, 0xF6)
GREEN = PptxRGB(0x22, 0xC5, 0x5E)
DARK = PptxRGB(0x0A, 0x0E, 0x14)
DARK_CARD = PptxRGB(0x11, 0x18, 0x27)
TEXT_LIGHT = PptxRGB(0xF1, 0xF5, 0xF9)
MUTED = PptxRGB(0x94, 0xA3, 0xB8)
WHITE = PptxRGB(0xFF, 0xFF, 0xFF)

DOCX_BLUE = RGBColor(0x3B, 0x82, 0xF6)
DOCX_INDIGO = RGBColor(0x63, 0x66, 0xF1)
DOCX_GREEN = RGBColor(0x22, 0xC5, 0x5E)
DOCX_DARK = RGBColor(0x11, 0x18, 0x27)
DOCX_MUTED = RGBColor(0x94, 0xA3, 0xB8)

SLIDES = [
    {
        "title": "Yoqubkhoja Hub",
        "subtitle": "Маркази рақамии лоиҳаҳо ва идоракунии ташкилот",
        "bullets": ["Семинари муаррифии барнома ва лоиҳа", "yoqubkhoja.tj"],
        "speaker": (
            "Муҳтарам ҳамкорон ва меҳмонони гиромӣ! Имрӯз ман барнома ва лоиҳаи "
            "Yoqubkhoja Hub-ро муаррифӣ мекунам — платформае, ки лоиҳаҳо, ташкилотҳо "
            "ва корҳои маъмуриро дар як ҷои ягона ҷамъ месозад."
        ),
    },
    {
        "title": "Чаро ин лоиҳа лозим шуд?",
        "bullets": [
            "Лоиҳаҳо дар ҷойҳои гуногун — GitHub, файлҳо, Excel, дафтарҳо",
            "Маълумоти кадрӣ ва молиявӣ пароканда ва дастӣ ҳисоб мешавад",
            "Ҳуҷҷатҳои меҳнатӣ, буҷет, табел — вақти зиёд мегирад",
            "4 забон — системаҳои кӯҳна инро дастгирӣ намекунанд",
        ],
        "speaker": (
            "Оё шумо ҳам вақте буд, ки барои як ҳисобот соатҳо ҷустуҷӯ кардед? "
            "Yoqubkhoja Hub барои ҳалли ҳамин мушкилот сохта шудааст."
        ),
    },
    {
        "title": "Yoqubkhoja Hub чист?",
        "image": "icon_tech",
        "bullets": [
            "Портали шахсии вебӣ — маркази лоиҳаҳо",
            "Воридшавии бехатар",
            "4 забон: тоҷикӣ, русӣ, англисӣ, ӯзбекӣ",
            "Модулҳои идоракунии ташкилот",
            "Сайт: https://yoqubkhoja.tj",
        ],
        "speaker": (
            "Ин на танҳо сайти шахсӣ аст — ин хаби рақамии кор аст: лоиҳаҳо, "
            "ташкилотҳо, кадрҳо, молия — ҳама дар як экран."
        ),
    },
    {
        "title": "Архитектураи умумӣ",
        "image": "diagram_arch",
        "bullets": [
            "Next.js 15 — зуд, муосир, барои production",
            "4 забон: тоҷикӣ · русӣ · англисӣ · ӯзбекӣ",
            "Амният: NextAuth, рамзи ҳашшуда",
            "Хостинг: Netlify + домени .tj",
        ],
        "speaker": (
            "Технологияҳои муосир, аммо мақсад содда аст: кор осонтар, "
            "хатогиҳо камтар, натиҷа зудтар."
        ),
    },
    {
        "title": "Модул 1: Маркази лоиҳаҳо",
        "image": "icon_projects",
        "bullets": [
            "Илова, таҳрир ва нест кардани лоиҳаҳо",
            "Ҷустуҷӯ, филтр, омори «дар кор / тайёр / нав»",
            "Категорияҳо: сомона, барнома, асбоб",
            "Пайванд ба GitHub",
        ],
        "speaker": (
            "Ҳар лоиҳа — як корт. Ҳамаи корҳои шумо дар як дашборд, "
            "на дар ёд ва на дар файлҳои гумшуда."
        ),
    },
    {
        "title": "Модул 2: Идоракунии ташкилотҳо",
        "image": "icon_org",
        "bullets": [
            "Реестри ташкилотҳо бо РМА",
            "Интегратсия бо andoz.tj",
            "Менюи фаъолият мувофиқи КПБ (cfs.tj)",
            "Ветеринария, фитосанитария, хизматрасониҳо",
        ],
        "speaker": (
            "Ташкилотро бо РМА меёбед — ном, суроға, маълумоти асосӣ "
            "худкор пур мешавад."
        ),
    },
    {
        "title": "Модул 3: Кормандон ва кадрҳо",
        "image": "icon_staff",
        "bullets": [
            "Омор — штат, бақайдгирифташуда, холигиҳо",
            "Басти вазифаҳо — ҳисоби автоматӣ",
            "Холигиҳо, бақайдгирӣ, табел",
            "Чоп ва экспорт",
        ],
        "speaker": "Кадрҳо — қалби ҳар ташкилот. Hub ин қалбро рақамӣ месозад.",
    },
    {
        "title": "Модул 4: Бухгалтерия ва молия",
        "image": "icon_finance",
        "bullets": [
            "Буҷет — нақша, иҷро, квартал",
            "Китоби музди меҳнат",
            "Пардохт ба бонк — Excel",
            "Рухсатии меҳнатӣ, ҳомиладорӣ, корношоямӣ",
        ],
        "speaker": (
            "Ҳуҷҷате, ки пештар соатҳо мегирифт, ҳоло дар чанд дақиқа — "
            "бо формулаҳои дуруст ва ҳуқуқии Тоҷикистон."
        ),
    },
    {
        "title": "Мувофиқати ҳуқуқӣ",
        "bullets": [
            "Корношоямӣ: музди вазифавӣ ÷ рӯзҳои меъёр × рӯзҳои беморӣ",
            "60% / 70% / 100% мувофиқи Қонуни суғурта",
            "Сақфи ҳадди ақал ва 2× музди вазифавӣ",
            "Кодекси меҳнат",
        ],
        "speaker": "Барнома на танҳо ҳисоб мекунад — қонунро эҳтиром мекунад.",
    },
    {
        "title": "Чандзабонӣ ва дастрасӣ",
        "bullets": ["🇹🇯 Тоҷикӣ", "🇷🇺 Русӣ", "🇬🇧 Англисӣ", "🇺🇿 Ӯзбекӣ"],
        "speaker": "Корманд бо забони дӯстдоштаи худ кор мекунад.",
    },
    {
        "title": "Экспорт ва ҳуҷҷатҳо",
        "bullets": ["Excel — бонк, табел", "PDF — ҳуҷҷатҳои расмӣ", "CSV — кормандон"],
        "speaker": "Маълумот ба бонк, ҳисобот ё чоп — як клик.",
    },
    {
        "title": "Намунаи амалӣ",
        "bullets": [
            "Маркази таъминоти бехатарии озуқаворӣ (Ҷ. Расулов)",
            "20 воҳиди корӣ, буҷети маҳаллӣ",
            "Хизматрасониҳои ветеринарӣ",
        ],
        "speaker": "Ин на танҳо концепт — лоиҳаи амалӣ аст.",
    },
    {
        "title": "Натиҷаҳо",
        "image": "diagram_results",
        "bullets": [
            "Лоиҳаҳо → як ҷо",
            "Ҳисоботи дастӣ → автоматӣ",
            "Excel → генератсияи онлайн",
            "1 забон → 4 забон",
        ],
        "speaker": "Вақт, хатогӣ ва стресс кам — сифат ва назорат зиёд.",
    },
    {
        "title": "Нақшаи оянда",
        "image": "icon_future",
        "bullets": [
            "Модулҳо барои ташкилотҳои дигар",
            "Интегратсияи давлатӣ",
            "Барномаи мобилӣ, API, open source",
        ],
        "speaker": "Hub метавонад платформаи миллӣ барои соҳаи озуқаворӣ бошад.",
    },
    {
        "title": "Хулоса",
        "bullets": [
            "Роҳи гузариш ба идоракунии рақамӣ",
            "yoqubkhoja.tj",
            "github.com/yoqubkhoja1988/my-yoqubkhoja-tj",
        ],
        "speaker": "Ин на танҳо барнома, балки тағйири усули кор аст.",
    },
    {
        "title": "Саволҳо?",
        "subtitle": "Ташаккур барои диққат!",
        "bullets": ["yoqubkhoja.tj", "github.com/yoqubkhoja1988"],
        "speaker": (
            "Агар савол дошта бошед — бо хушҳолӣ ҷавоб медиҳам. "
            "Ташаккур барои диққат!"
        ),
        "closing": True,
    },
]

ASSETS: dict[str, Path] = {}


def set_paragraph_style(paragraph, size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
    paragraph.font.size = PptxPt(size)
    paragraph.font.bold = bold
    if color:
        paragraph.font.color.rgb = color
    paragraph.alignment = align


def add_full_bleed_image(slide, image_path: Path, slide_width, slide_height):
    slide.shapes.add_picture(str(image_path), 0, 0, width=slide_width, height=slide_height)


def add_brand_footer(slide, slide_num: int | None = None, dark=False):
    logo_path = ASSETS.get("logo_sm")
    if logo_path and logo_path.exists():
        slide.shapes.add_picture(str(logo_path), PptxInches(0.35), PptxInches(6.85), height=PptxInches(0.42))

    footer = slide.shapes.add_textbox(PptxInches(1.0), PptxInches(6.9), PptxInches(7.5), PptxInches(0.35))
    fp = footer.text_frame.paragraphs[0]
    fp.text = "Yoqubkhoja Hub  ·  yoqubkhoja.tj  ·  2026"
    set_paragraph_style(fp, size=9, color=TEXT_LIGHT if dark else MUTED, align=PP_ALIGN.LEFT)

    if slide_num is not None:
        num = slide.shapes.add_textbox(PptxInches(9.1), PptxInches(6.9), PptxInches(0.6), PptxInches(0.35))
        np = num.text_frame.paragraphs[0]
        np.text = str(slide_num)
        set_paragraph_style(np, size=9, color=TEXT_LIGHT if dark else MUTED, align=PP_ALIGN.RIGHT)


def add_gradient_bar(slide, colors=(BLUE, INDIGO, PURPLE, GREEN)):
    width = PptxInches(10) / len(colors)
    for i, color in enumerate(colors):
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, width * i, 0, width, PptxInches(0.14))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()


def add_title_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = ASSETS.get("bg_dark")
    if bg and bg.exists():
        add_full_bleed_image(slide, bg, prs.slide_width, prs.slide_height)

    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, PptxInches(10), PptxInches(7.5))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = DARK
    overlay.fill.transparency = 0.25
    overlay.line.fill.background()

    logo = ASSETS.get("logo")
    if logo and logo.exists():
        slide.shapes.add_picture(str(logo), PptxInches(4.05), PptxInches(0.9), height=PptxInches(1.35))

    box = slide.shapes.add_textbox(PptxInches(0.6), PptxInches(2.5), PptxInches(8.8), PptxInches(2.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Yoqubkhoja Hub"
    set_paragraph_style(p, size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    sub = tf.add_paragraph()
    sub.text = "Семинари муаррифии барнома ва лоиҳа"
    set_paragraph_style(sub, size=22, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

    sub2 = tf.add_paragraph()
    sub2.text = "Маркази рақамии лоиҳаҳо ва идоракунии ташкилот"
    set_paragraph_style(sub2, size=16, color=MUTED, align=PP_ALIGN.CENTER)

    tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptxInches(3.1), PptxInches(5.0), PptxInches(3.8), PptxInches(0.55))
    tag.fill.solid()
    tag.fill.fore_color.rgb = BLUE
    tag.line.fill.background()
    tag_tf = tag.text_frame
    tag_tf.paragraphs[0].text = "yoqubkhoja.tj"
    set_paragraph_style(tag_tf.paragraphs[0], size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tag_tf.vertical_anchor = 1

    add_brand_footer(slide, dark=True)


def add_content_slide(prs: Presentation, index: int, data: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    closing = data.get("closing", False)

    if closing:
        bg = ASSETS.get("bg_dark")
        if bg and bg.exists():
            add_full_bleed_image(slide, bg, prs.slide_width, prs.slide_height)
        overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, PptxInches(10), PptxInches(7.5))
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = DARK
        overlay.fill.transparency = 0.3
        overlay.line.fill.background()
        title_color = WHITE
        body_color = TEXT_LIGHT
        muted_color = MUTED
    else:
        bg = ASSETS.get("bg_light")
        if bg and bg.exists():
            add_full_bleed_image(slide, bg, prs.slide_width, prs.slide_height)
        add_gradient_bar(slide)
        title_color = BLUE
        body_color = DARK_CARD
        muted_color = MUTED

    has_image = bool(data.get("image")) and ASSETS.get(data["image"])
    text_width = PptxInches(4.8) if has_image else PptxInches(8.6)

    logo = ASSETS.get("logo_sm")
    if logo and logo.exists():
        slide.shapes.add_picture(str(logo), PptxInches(8.95), PptxInches(0.22), height=PptxInches(0.45))

    title_box = slide.shapes.add_textbox(PptxInches(0.55), PptxInches(0.35), text_width, PptxInches(0.95))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_p = title_tf.paragraphs[0]
    title_p.text = data["title"]
    set_paragraph_style(title_p, size=28, bold=True, color=title_color)

    if data.get("subtitle"):
        sub_p = title_tf.add_paragraph()
        sub_p.text = data["subtitle"]
        set_paragraph_style(sub_p, size=16, color=muted_color)

    body_top = 1.35 if not data.get("subtitle") else 1.65
    body_box = slide.shapes.add_textbox(PptxInches(0.7), PptxInches(body_top), text_width, PptxInches(3.5))
    body_tf = body_box.text_frame
    body_tf.word_wrap = True

    for i, bullet in enumerate(data.get("bullets", [])):
        p = body_tf.paragraphs[0] if i == 0 else body_tf.add_paragraph()
        p.text = f"• {bullet}"
        set_paragraph_style(p, size=16, color=body_color)
        p.space_after = PptxPt(6)

    if has_image:
        image_path = ASSETS[data["image"]]
        slide.shapes.add_picture(
            str(image_path),
            PptxInches(5.35),
            PptxInches(1.2),
            width=PptxInches(4.2),
        )

    if data.get("speaker"):
        speaker_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            PptxInches(0.55),
            PptxInches(5.35),
            PptxInches(8.9),
            PptxInches(1.15),
        )
        speaker_box.fill.solid()
        if closing:
            speaker_box.fill.fore_color.rgb = BLUE_DARK
            speaker_box.fill.transparency = 0.05
            speaker_box.line.color.rgb = INDIGO
            speaker_text_color = TEXT_LIGHT
        else:
            speaker_box.fill.fore_color.rgb = WHITE
            speaker_box.fill.transparency = 0.05
            speaker_box.line.color.rgb = BLUE
            speaker_text_color = DARK_CARD
        speaker_tf = speaker_box.text_frame
        speaker_tf.word_wrap = True
        speaker_tf.margin_left = PptxInches(0.15)
        speaker_tf.margin_right = PptxInches(0.15)
        speaker_tf.margin_top = PptxInches(0.1)
        sp = speaker_tf.paragraphs[0]
        sp.text = f"Суханрон: {data['speaker']}"
        set_paragraph_style(sp, size=11, color=speaker_text_color)

    if closing and logo and logo.exists():
        slide.shapes.add_picture(str(logo), PptxInches(4.05), PptxInches(2.0), height=PptxInches(1.0))

    add_brand_footer(slide, index, dark=closing)


def build_pptx():
    prs = Presentation()
    prs.slide_width = PptxInches(10)
    prs.slide_height = PptxInches(7.5)

    add_title_slide(prs)
    for i, slide_data in enumerate(SLIDES, start=1):
        add_content_slide(prs, i, slide_data)

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(PPTX_PATH)
    return PPTX_PATH


def set_docx_heading_color(paragraph, color: RGBColor):
    for run in paragraph.runs:
        run.font.color.rgb = color


def add_horizontal_line(doc):
    p = doc.add_paragraph()
    p_pr = p._p.get_or_add_pPr()
    p_bdr = OxmlElement("w:pBdr")
    bottom = OxmlElement("w:bottom")
    bottom.set(qn("w:val"), "single")
    bottom.set(qn("w:sz"), "12")
    bottom.set(qn("w:space"), "1")
    bottom.set(qn("w:color"), "3B82F6")
    p_bdr.append(bottom)
    p_pr.append(p_bdr)


def build_docx():
    doc = Document()
    section = doc.sections[0]
    section.top_margin = Inches(0.9)
    section.bottom_margin = Inches(0.9)
    section.left_margin = Inches(1)
    section.right_margin = Inches(1)

    logo = ASSETS.get("logo")
    if logo and logo.exists():
        doc.add_picture(str(logo), width=Inches(1.2))
        last = doc.paragraphs[-1]
        last.alignment = WD_ALIGN_PARAGRAPH.CENTER

    title = doc.add_heading("Yoqubkhoja Hub", level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    set_docx_heading_color(title, DOCX_BLUE)

    subtitle = doc.add_paragraph("Семинари муаррифии барнома ва лоиҳа")
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle.runs[0].font.size = Pt(14)
    subtitle.runs[0].font.color.rgb = DOCX_INDIGO
    subtitle.runs[0].bold = True

    tag = doc.add_paragraph("Маркази рақамии лоиҳаҳо ва идоракунии ташкилот")
    tag.alignment = WD_ALIGN_PARAGRAPH.CENTER
    tag.runs[0].italic = True
    tag.runs[0].font.color.rgb = DOCX_MUTED

    add_horizontal_line(doc)

    links = doc.add_paragraph()
    links.alignment = WD_ALIGN_PARAGRAPH.CENTER
    links.add_run("Сайт: ").bold = True
    links.add_run("https://yoqubkhoja.tj\n")
    links.add_run("GitHub: ").bold = True
    links.add_run("https://github.com/yoqubkhoja1988/my-yoqubkhoja-tj")

    hero = ASSETS.get("bg_dark")
    if hero and hero.exists():
        doc.add_picture(str(hero), width=Inches(6))
        doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_page_break()

    h = doc.add_heading("Мундариҷа", level=1)
    set_docx_heading_color(h, DOCX_BLUE)
    for i, slide in enumerate(SLIDES, start=1):
        doc.add_paragraph(f"{i}. {slide['title']}", style="List Number")

    doc.add_page_break()

    h = doc.add_heading("Маслиҳат барои суханронӣ (5–7 дақиқа)", level=1)
    set_docx_heading_color(h, DOCX_BLUE)
    timing = [
        ("Слайдҳои 1–3", "1 дақ — муаррифӣ"),
        ("Слайдҳои 4–5", "1 дақ — техникаву лоиҳаҳо"),
        ("Слайдҳои 6–9", "2–3 дақ — ташкилот, кадрҳо, молия"),
        ("Слайдҳои 10–16", "1–2 дақ — натиҷа, оянда, хулоса"),
    ]
    for label, value in timing:
        p = doc.add_paragraph(style="List Bullet")
        p.add_run(f"{label}: ").bold = True
        p.add_run(value)

    doc.add_page_break()

    for i, slide in enumerate(SLIDES, start=1):
        h = doc.add_heading(f"Слайд {i} — {slide['title']}", level=1)
        set_docx_heading_color(h, DOCX_INDIGO if i % 2 else DOCX_BLUE)

        image_key = slide.get("image")
        if image_key and ASSETS.get(image_key) and ASSETS[image_key].exists():
            doc.add_picture(str(ASSETS[image_key]), width=Inches(4.8))

        if slide.get("subtitle"):
            p = doc.add_paragraph(slide["subtitle"])
            p.runs[0].italic = True
            p.runs[0].font.color.rgb = DOCX_MUTED

        if slide.get("bullets"):
            doc.add_paragraph("Матни слайд:", style="Heading 3")
            for bullet in slide["bullets"]:
                bp = doc.add_paragraph(bullet, style="List Bullet")
                for run in bp.runs:
                    run.font.color.rgb = DOCX_DARK

        if slide.get("speaker"):
            doc.add_paragraph("Суханрон:", style="Heading 3")
            speaker = doc.add_paragraph(slide["speaker"])
            speaker.runs[0].italic = True
            speaker.runs[0].font.color.rgb = DOCX_GREEN

        if i < len(SLIDES):
            add_horizontal_line(doc)

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    doc.save(DOCX_PATH)
    return DOCX_PATH


def main():
    global ASSETS
    ASSETS = build_assets()
    pptx = build_pptx()
    docx = build_docx()
    print(f"Assets: {ASSETS['logo'].parent}")
    print(f"Created: {pptx}")
    print(f"Created: {docx}")


if __name__ == "__main__":
    main()
